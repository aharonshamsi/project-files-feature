import json
import os
import re
import io
import hashlib
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from src.models.document_models import Metadata, ContentBlock, DocumentModel, TableData, ImageData
from src.parsers.utils import  is_solid_color_image

# =========================================================
# Constants
HEADING_TYPE = "heading"
PARAGRAPH_TYPE = "paragraph"
TABLE_TYPE = "table"
URL_TYPE = "url"

DEFAULT_FONT_SIZE = 14.0
FONT_SIZE_DELTA = 5
TOP_HEADING_RATIO = 0.2
MAX_HEADING_WORDS = 10
EMPHASIZED_RATIO = 0.6
MINI_WORDS = 40 
MIN_IMAGE_SIZE_BYTES = 2*1024 # Minimum image size threshold (2KB)

# =========================================================
def extract_pptx_metadata(prs) -> Metadata:
    """
    Extracts core properties and metadata from the PowerPoint presentation.

    Args:
        prs: The python-pptx Presentation object.

    Returns:
        Metadata: A model containing presentation title, author, and creation date.
    """
    props = prs.core_properties

    metadata = {
        "title": props.title or "",
        "author": props.author or "",
        "creation_date": str(props.created) if props.created else None,
    }

    return Metadata(**metadata)

# =========================================================
def normalize_text(text):
    """
    Cleans up raw PowerPoint text, replacing weird control characters with newlines.

    Args:
        text (str): The raw text to format.

    Returns:
        str: The normalized, stripped string.
    """
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

# =========================================================
def extract_urls(text):
    """
    Finds HTTP/HTTPS URLs within a block of text using regex.

    Args:
        text (str): The text block to scan.

    Returns:
        list: A list of matched URL strings.
    """
    return re.findall(r'https?://[^\s)]+', text)

# =========================================================
def is_auto_slide_number(shape):
    """
    Determines if a PowerPoint shape is an auto-generated slide number placeholder.

    Args:
        shape: The shape object from a slide.

    Returns:
        bool: True if it's a slide number placeholder, False otherwise.
    """
    if not shape.is_placeholder:
        return False
    try:
        return shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER
    except KeyError:
        return False

# =========================================================
def get_slide_body_size(slide):
    """
    Calculates the most frequent font size used in the text frames of a single slide.

    Args:
        slide: A python-pptx slide object.

    Returns:
        float: The determined body font size (mode) or DEFAULT_FONT_SIZE.
    """
    font_counts = {}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    size = round(run.font.size.pt, 1)
                    font_counts[size] = font_counts.get(size, 0) + 1

    return max(font_counts, key=font_counts.get) if font_counts else DEFAULT_FONT_SIZE

# =========================================================
def shape_is_heading(shape, slide_height, body_font_size):
    """
    Determines if a given text shape acts as a heading based on placeholder type, 
    position, word count, and text emphasis.

    Args:
        shape: The slide shape to evaluate.
        slide_height (int): The total height of the slide.
        body_font_size (float): The calculated baseline font size for the slide.

    Returns:
        bool: True if the shape qualifies as a heading, False otherwise.
    """
    if not shape.has_text_frame:
        return False

    text = normalize_text(shape.text_frame.text)
    if not text:
        return False

    words = text.split()

    # Placeholder title always wins
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return True
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return False

    shape_bottom = shape.top + shape.height

    if (
        shape_bottom < slide_height * TOP_HEADING_RATIO and
        len(words) <= MAX_HEADING_WORDS
    ):
        emphasized_runs = 0
        total_runs = 0

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                total_runs += 1
                if run.font.bold or (
                    run.font.size and run.font.size.pt > body_font_size + FONT_SIZE_DELTA
                ):
                    emphasized_runs += 1

        if total_runs and emphasized_runs / total_runs >= EMPHASIZED_RATIO:
            return True

    return False


# =========================================================
def extract_pptx_table(table, count_words):
    """
    Extracts text contents from a PowerPoint table and updates the word count.

    Args:
        table: The python-pptx Table object.
        count_words (int): Current accumulated word count.

    Returns:
        tuple: A dictionary with 'headers' and 'rows', and the updated word count.
    """
    extracted = []

    for row in table.rows:
        extracted_row = []
        for cell in row.cells:
            text = normalize_text(cell.text_frame.text)
            extracted_row.append(text)
            count_words += len(text.split())
        extracted.append(extracted_row)

    if extracted:
        table_data = {
            "headers": extracted[0],
            "rows": extracted[1:]
        }
    else:
        table_data = {
            "headers": [],
            "rows": []
        }

    return table_data, count_words


# =========================================================
def extract_pptx_file_to_model(file_stream: io.BytesIO, image_output_dir: str) -> tuple[int, DocumentModel]:
    """
    Main parser for PPTX files. Iterates through slides, shapes, and tables,
    categorizing text and extracting images into a structured DocumentModel.

    Args:
        file_stream (io.BytesIO): Raw binary stream of the PPTX file.
        image_output_dir (str): Directory where extracted images are saved.

    Returns:
        tuple[int, DocumentModel]: Total word count and the structured DocumentModel.

    Raises:
        ValueError: If total words are below the MINI_WORDS threshold.
    """
    file_stream.seek(0) 
    prs = Presentation(file_stream)

    slide_height = prs.slide_height
    total_word_count = 0
    block_id_counter = 1

    metadata = extract_pptx_metadata(prs)
    block_list = [] 

    try:
        for slide_index, slide in enumerate(prs.slides):
            body_font_size = get_slide_body_size(slide)
            shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))

            for shape in shapes:
                if is_auto_slide_number(shape):
                    continue


                # ================= IMAGE =================
                if shape.shape_type == 13:  # Picture
                    image_paths = extract_and_save_pptx_images(shape, block_id_counter, image_output_dir)

                    for path in image_paths:
                        block_list.append(ContentBlock(
                            block_id=block_id_counter,
                            type="image",
                            image_data=ImageData(image_path=path)
                        ))
                        block_id_counter += 1
                    continue


                # ================= TABLE =================
                if shape.has_table:
                    table_obj, total_word_count = extract_pptx_table(shape.table, total_word_count)

                    block_list.append(ContentBlock(
                        block_id=block_id_counter,
                        type=TABLE_TYPE,
                        table_data=TableData(**table_obj)
                    ))
                    block_id_counter += 1
                    continue

                # ================= TEXT =================
                if not shape.has_text_frame:
                    continue

                element_type = (
                    HEADING_TYPE
                    if shape_is_heading(shape, slide_height, body_font_size)
                    else PARAGRAPH_TYPE
                )

                collected_text = []

                for paragraph in shape.text_frame.paragraphs:
                    text = normalize_text(paragraph.text)
                    if text:
                        collected_text.append(text)

                if not collected_text:
                    continue

                full_text = "\n".join(collected_text)

                urls = extract_urls(full_text)
                clean_text = full_text
                for url in urls:
                    clean_text = clean_text.replace(url, "").strip()

                if clean_text:
                    total_word_count += len(clean_text.split())
                    block_list.append(ContentBlock(
                        block_id=block_id_counter,
                        type=element_type,
                        text=clean_text
                    ))
                    block_id_counter += 1

                for url in urls:
                    total_word_count += 1
                    block_list.append(ContentBlock(
                        block_id=block_id_counter,
                        type=URL_TYPE,
                        text=url
                    ))
                    block_id_counter += 1

        if total_word_count < MINI_WORDS:
            raise ValueError(f"Content too short: {total_word_count} words.")

        final_document = DocumentModel(
            metadata=metadata,
            content_blocks=block_list
        )

        return total_word_count, final_document
    
    except Exception as e:
        print(f"Error: {e}")
        raise


#===============================================================================
def extract_and_save_pptx_images(shape, block_id, image_output_dir: str):
    """
    Extracts binary data from a PowerPoint picture shape, filters out solids/small 
    images, and saves it locally using a SHA256 filename.

    Args:
        shape: The python-pptx picture shape.
        block_id (int): The current identifier for logging purposes.
        image_output_dir (str): Destination directory for the saved image.

    Returns:
        list: A list containing the absolute path string to the saved image (or empty if skipped).
    """
    image_paths = []

    if not shape.shape_type == 13:  # 13 == MSO_SHAPE_TYPE.PICTURE
        return image_paths

    try:
        image = shape.image
        image_bytes = image.blob

        # Skip saving if the image is smaller than the threshold
        if len(image_bytes) < MIN_IMAGE_SIZE_BYTES:
            return image_paths # Returns an empty list

        if is_solid_color_image(image_bytes):
            print(f"Skipping solid color image found in PowerPoint block {block_id}")
            return image_paths # Returns an empty list

        extension = image.ext
        # Generate unique hash from the image bytes
        image_hash = hashlib.sha256(image_bytes).hexdigest()
        image_filename = f"{image_hash}.{extension}"
        full_path = os.path.join(image_output_dir, image_filename)

        # Check if the image already exists on disk before saving
        if not os.path.exists(full_path):
            with open(full_path, "wb") as f:
                f.write(image_bytes)

        image_paths.append(full_path)

    except Exception as e:
        raise RuntimeError(f"Image extraction failed for block {block_id}: {e}") from e

    return image_paths